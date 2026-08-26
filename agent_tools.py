"""
Agent Tools — Les 4 outils de l'agent (V2).
V1: list_dir, read_file (lecture seule)
V2: + copy_file, write_file (écriture contrôlée)
Chaque outil passe par PathGuard avant toute opération I/O.
"""

import os
import shutil
import io
from path_guard import PathGuard

try:
    import pandas as pd
    import docx
    HAS_OFFICE_LIBS = True
except ImportError:
    HAS_OFFICE_LIBS = False

try:
    import pypdf
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    import pptx
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# Max characters to return from a file read (protects LLM context window)
MAX_CONTENT_LENGTH = 4000


# ─── TOOL 1: list_dir ───────────────────────────────────────

def list_dir(path: str, guard: PathGuard) -> dict:
    """List files and subdirectories at the given path."""
    allowed, reason = guard.validate(path, action="list_dir")
    if not allowed:
        return {"success": False, "error": reason}

    resolved = guard.resolve(path)

    if not os.path.isdir(resolved):
        return {"success": False, "error": f"Ce n'est pas un dossier : {resolved}"}

    entries = []
    try:
        for item in sorted(os.listdir(resolved)):
            full_path = os.path.join(resolved, item)
            if os.path.isdir(full_path):
                entries.append(f"📁 {item}/")
            else:
                size = os.path.getsize(full_path)
                if size < 1024:
                    size_str = f"{size} B"
                elif size < 1024 * 1024:
                    size_str = f"{size / 1024:.1f} KB"
                else:
                    size_str = f"{size / (1024 * 1024):.1f} MB"
                entries.append(f"📄 {item} ({size_str})")
    except PermissionError:
        return {"success": False, "error": f"Permission refusée : {resolved}"}

    return {
        "success": True,
        "path": resolved,
        "count": len(entries),
        "entries": entries,
    }


# ─── TOOL 2: read_file ──────────────────────────────────────

def read_file(path: str, guard: PathGuard) -> dict:
    """Read the text content of a file (truncated if too long)."""
    allowed, reason = guard.validate(path, action="read_file")
    if not allowed:
        return {"success": False, "error": reason}

    resolved = guard.resolve(path)

    if not os.path.isfile(resolved):
        return {"success": False, "error": f"Ce n'est pas un fichier : {resolved}"}

    ext = os.path.splitext(resolved)[1].lower()

    try:
        if ext == ".docx":
            if not HAS_OFFICE_LIBS:
                return {"success": False, "error": "Librairies manquantes. Lancez: pip install python-docx"}
            doc = docx.Document(resolved)
            content = "\n".join([p.text for p in doc.paragraphs])
        elif ext == ".xlsx":
            if not HAS_OFFICE_LIBS:
                return {"success": False, "error": "Librairies manquantes. Lancez: pip install pandas openpyxl"}
            df = pd.read_excel(resolved)
            content = df.to_csv(index=False)
        elif ext == ".pdf":
            if not HAS_PDF:
                return {"success": False, "error": "Librairie manquante. Lancez: pip install pypdf"}
            reader = pypdf.PdfReader(resolved)
            pages = []
            for i, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                if text.strip():
                    pages.append(f"[Page {i+1}]\n{text}")
            content = "\n\n".join(pages) if pages else "(PDF vide ou sans texte extractible)"
        elif ext == ".pptx":
            if not HAS_PPTX:
                return {"success": False, "error": "Librairie manquante. Lancez: pip install python-pptx"}
            prs = pptx.Presentation(resolved)
            slides = []
            for i, slide in enumerate(prs.slides):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                texts.append(para.text)
                if texts:
                    slides.append(f"[Slide {i+1}]\n" + "\n".join(texts))
            content = "\n\n".join(slides) if slides else "(Présentation vide ou sans texte)"
        else:
            with open(resolved, "r", encoding="utf-8", errors="replace") as f:
                content = f.read(MAX_CONTENT_LENGTH + 1)

        truncated = len(content) > MAX_CONTENT_LENGTH
        if truncated:
            content = content[:MAX_CONTENT_LENGTH]

        return {
            "success": True,
            "path": resolved,
            "content": content,
            "truncated": truncated,
        }
    except Exception as e:
        return {"success": False, "error": f"Erreur de lecture : {str(e)}"}


# ─── TOOL 3: copy_file ──────────────────────────────────────

def copy_file(args: str, guard: PathGuard) -> dict:
    """
    Copy a file from source to destination.
    Args format: "source | destination"
    """
    parts = args.split("|", 1)
    if len(parts) < 2:
        return {"success": False, "error": "Format requis : copy_file | source | destination"}

    source = parts[0].strip()
    destination = parts[1].strip()

    # Validate both paths via PathGuard
    allowed, reason = guard.validate_copy(source, destination)
    if not allowed:
        return {"success": False, "error": reason}

    resolved_src = guard.resolve(source)
    resolved_dst = guard.resolve(destination)

    if not os.path.isfile(resolved_src):
        return {"success": False, "error": f"La source n'est pas un fichier : {resolved_src}"}

    # Backup destination if it already exists (Rule 7)
    backup_path = guard.backup_if_exists(destination)

    try:
        # Create parent directories if needed
        os.makedirs(os.path.dirname(resolved_dst), exist_ok=True)
        shutil.copy2(resolved_src, resolved_dst)

        result = {
            "success": True,
            "source": resolved_src,
            "destination": resolved_dst,
            "message": f"Fichier copié avec succès.",
        }
        if backup_path:
            result["backup"] = backup_path
            result["message"] += f" Backup de l'ancien fichier : {backup_path}"
        return result

    except Exception as e:
        return {"success": False, "error": f"Erreur de copie : {str(e)}"}


# ─── TOOL 4: write_file ─────────────────────────────────────

def write_file(args: str, guard: PathGuard) -> dict:
    """
    Write content to a file (create or overwrite).
    Args format: "path | content"
    """
    parts = args.split("|", 1)
    if len(parts) < 2:
        return {"success": False, "error": "Format requis : write_file | chemin | contenu"}

    path = parts[0].strip()
    content = parts[1].strip()

    # Validate path via PathGuard
    allowed, reason = guard.validate(path, action="write_file")
    if not allowed:
        return {"success": False, "error": reason}

    # Validate content size (Rule 8)
    ok, size_reason = guard.validate_write_content(content)
    if not ok:
        return {"success": False, "error": size_reason}

    resolved = guard.resolve(path)
    ext = os.path.splitext(resolved)[1].lower()

    # Backup if file already exists (Rule 7)
    backup_path = guard.backup_if_exists(path)

    try:
        # Create parent directories if needed
        os.makedirs(os.path.dirname(resolved), exist_ok=True)
        
        if ext == ".docx":
            if not HAS_OFFICE_LIBS:
                return {"success": False, "error": "Librairies manquantes. Lancez: pip install pandas openpyxl python-docx"}
            doc = docx.Document()
            doc.add_paragraph(content)
            doc.save(resolved)
        elif ext == ".xlsx":
            if not HAS_OFFICE_LIBS:
                return {"success": False, "error": "Librairies manquantes. Lancez: pip install pandas openpyxl python-docx"}
            df = pd.read_csv(io.StringIO(content))
            df.to_excel(resolved, index=False)
        else:
            with open(resolved, "w", encoding="utf-8") as f:
                f.write(content)

        result = {
            "success": True,
            "path": resolved,
            "chars_written": len(content),
            "message": f"Fichier écrit avec succès ({len(content)} caractères).",
        }
        if backup_path:
            result["backup"] = backup_path
            result["message"] += f" Backup de l'ancien fichier : {backup_path}"
        return result

    except Exception as e:
        return {"success": False, "error": f"Erreur d'écriture : {str(e)}"}


# ─── TOOL 5: move_file ──────────────────────────────────────

def move_file(args: str, guard: PathGuard) -> dict:
    """
    Move (or rename) a file from source to destination.
    Args format: "source | destination"
    """
    parts = args.split("|", 1)
    if len(parts) < 2:
        return {"success": False, "error": "Format requis : move_file | source | destination"}

    source = parts[0].strip()
    destination = parts[1].strip()

    # Validate both paths via PathGuard
    allowed, reason = guard.validate_copy(source, destination)
    if not allowed:
        return {"success": False, "error": reason}

    resolved_src = guard.resolve(source)
    resolved_dst = guard.resolve(destination)

    if not os.path.isfile(resolved_src):
        return {"success": False, "error": f"La source n'est pas un fichier : {resolved_src}"}

    # Backup destination if it already exists (Rule 7)
    backup_path = guard.backup_if_exists(destination)

    try:
        # Create parent directories if needed
        os.makedirs(os.path.dirname(resolved_dst), exist_ok=True)
        shutil.move(resolved_src, resolved_dst)

        result = {
            "success": True,
            "source": resolved_src,
            "destination": resolved_dst,
            "message": f"Fichier déplacé avec succès.",
        }
        if backup_path:
            result["backup"] = backup_path
            result["message"] += f" Backup de l'ancien fichier : {backup_path}"
        return result

    except Exception as e:
        return {"success": False, "error": f"Erreur de déplacement : {str(e)}"}


# ─── TOOL 6: search_in_files ────────────────────────────────

def search_in_files(args: str, guard: PathGuard) -> dict:
    """
    Search for a text pattern inside files in a directory.
    Args format: "directory | pattern"
    """
    parts = args.split("|", 1)
    if len(parts) < 2:
        return {"success": False, "error": "Format requis : search_in_files | dossier | texte_a_chercher"}

    directory = parts[0].strip()
    pattern = parts[1].strip().lower()

    allowed, reason = guard.validate(directory, action="list_dir")
    if not allowed:
        return {"success": False, "error": reason}

    resolved_dir = guard.resolve(directory)
    if not os.path.isdir(resolved_dir):
        return {"success": False, "error": f"Dossier introuvable : {resolved_dir}"}

    results = []
    try:
        for root, _, files in os.walk(resolved_dir):
            for file in files:
                ext = os.path.splitext(file)[1].lower()
                # On limite la recherche aux fichiers textes pour le grep basique
                if ext in [".txt", ".md", ".csv", ".json", ".yaml", ".yml", ".py", ".log"]:
                    full_path = os.path.join(root, file)
                    try:
                        with open(full_path, "r", encoding="utf-8", errors="ignore") as f:
                            for i, line in enumerate(f):
                                if pattern in line.lower():
                                    rel_path = os.path.relpath(full_path, guard.root_dir)
                                    results.append(f"{rel_path}:{i+1}: {line.strip()[:150]}")
                                    if len(results) >= 50: # Limite pour éviter de surcharger le LLM
                                        results.append("... [Plus de 50 résultats, recherche tronquée]")
                                        break
                    except Exception:
                        pass
            if len(results) >= 50:
                break
    except Exception as e:
        return {"success": False, "error": f"Erreur de recherche : {str(e)}"}

    return {
        "success": True,
        "pattern": pattern,
        "count": len(results),
        "results": results,
    }


# ─── TOOL 7: delete_file ────────────────────────────────────

def delete_file(path: str, guard: PathGuard) -> dict:
    """Delete a file with automatic backup to a recycle bin (or .bak)."""
    allowed, reason = guard.validate(path, action="delete_file")
    if not allowed:
        return {"success": False, "error": reason}

    resolved = guard.resolve(path)
    if not os.path.isfile(resolved):
        return {"success": False, "error": f"Fichier introuvable : {resolved}"}

    # Backup before deletion
    backup_path = guard.backup_if_exists(path)

    try:
        os.remove(resolved)
        return {
            "success": True,
            "path": resolved,
            "message": f"Fichier supprimé. Backup sauvegardé dans : {backup_path}"
        }
    except Exception as e:
        return {"success": False, "error": f"Erreur de suppression : {str(e)}"}


# ─── TOOL 8: create_dir ─────────────────────────────────────

def create_dir(path: str, guard: PathGuard) -> dict:
    """Create a new directory."""
    allowed, reason = guard.validate(path, action="create_dir")
    if not allowed:
        return {"success": False, "error": reason}

    resolved = guard.resolve(path)
    if os.path.exists(resolved):
        return {"success": False, "error": f"Le chemin existe déjà : {resolved}"}

    try:
        os.makedirs(resolved)
        return {
            "success": True,
            "path": resolved,
            "message": "Dossier créé avec succès."
        }
    except Exception as e:
        return {"success": False, "error": f"Erreur de création de dossier : {str(e)}"}


# ─── TOOL 9: append_to_file ─────────────────────────────────

def append_to_file(args: str, guard: PathGuard) -> dict:
    """
    Append text to the end of an existing file.
    Args format: "path | content"
    """
    parts = args.split("|", 1)
    if len(parts) < 2:
        return {"success": False, "error": "Format requis : append_to_file | chemin | contenu"}

    path = parts[0].strip()
    content = parts[1] # On garde les espaces ici, c'est de l'ajout

    allowed, reason = guard.validate(path, action="append_to_file")
    if not allowed:
        return {"success": False, "error": reason}

    resolved = guard.resolve(path)
        
    ext = os.path.splitext(resolved)[1].lower()
    if ext in [".docx", ".xlsx", ".pdf", ".pptx"]:
        return {"success": False, "error": f"Impossible d'utiliser append_to_file sur des fichiers binaires ({ext}). Utilisez write_file."}

    # Backup before appending (only if it exists)
    backup_path = guard.backup_if_exists(path)

    try:
        with open(resolved, "a", encoding="utf-8") as f:
            if not content.startswith("\n"):
                f.write("\n")
            f.write(content)
        
        result = {
            "success": True,
            "path": resolved,
            "chars_added": len(content),
            "message": f"Texte ajouté avec succès ({len(content)} caractères)."
        }
        if backup_path:
            result["backup"] = backup_path
            result["message"] += f" Backup de sécurité créé avant l'ajout : {backup_path}"
        return result
    except Exception as e:
        return {"success": False, "error": f"Erreur d'ajout : {str(e)}"}


# ─── TOOL 10: file_info ─────────────────────────────────────
import time

def file_info(path: str, guard: PathGuard) -> dict:
    """Get metadata about a file or directory."""
    allowed, reason = guard.validate(path, action="list_dir")
    if not allowed:
        return {"success": False, "error": reason}

    resolved = guard.resolve(path)
    if not os.path.exists(resolved):
        return {"success": False, "error": f"Chemin introuvable : {resolved}"}

    stat = os.stat(resolved)
    is_dir = os.path.isdir(resolved)
    
    info = {
        "success": True,
        "path": resolved,
        "type": "directory" if is_dir else "file",
        "size_bytes": stat.st_size,
        "created": time.ctime(stat.st_ctime),
        "modified": time.ctime(stat.st_mtime),
    }
    if not is_dir:
        info["extension"] = os.path.splitext(resolved)[1].lower()
    return info


# ─── TOOL 11: tree ──────────────────────────────────────────

def tree(path: str, guard: PathGuard) -> dict:
    """Return the recursive directory structure as text."""
    allowed, reason = guard.validate(path, action="list_dir")
    if not allowed:
        return {"success": False, "error": reason}

    resolved = guard.resolve(path)
    if not os.path.isdir(resolved):
        return {"success": False, "error": f"Dossier introuvable : {resolved}"}

    def build_tree(dir_path, prefix="", depth=0, max_depth=3):
        if depth > max_depth:
            return [prefix + "└── ... (limite de profondeur atteinte)"]
        
        try:
            items = sorted(os.listdir(dir_path))
        except PermissionError:
            return [prefix + "└── [Accès refusé]"]
            
        lines = []
        for i, item in enumerate(items):
            is_last = (i == len(items) - 1)
            connector = "└── " if is_last else "├── "
            full_path = os.path.join(dir_path, item)
            
            if os.path.isdir(full_path):
                lines.append(prefix + connector + "📁 " + item)
                extension = "    " if is_last else "│   "
                lines.extend(build_tree(full_path, prefix + extension, depth + 1, max_depth))
            else:
                lines.append(prefix + connector + "📄 " + item)
        return lines

    tree_lines = [os.path.basename(resolved) or "."] + build_tree(resolved)
    tree_output = "\n".join(tree_lines)
    
    # Sécurité anti-spam pour le contexte du LLM
    if len(tree_output) > 10000:
        tree_output = tree_output[:10000] + "\n... [ARBORESCENCE TRONQUÉE (trop grande)]"

    return {
        "success": True,
        "path": resolved,
        "tree": tree_output
    }


# ─── Tool registry (used by agent_core) ────────────────────

TOOLS_DESCRIPTION = """
You have access to exactly 11 tools. Use them to help the user explore and manage files.

TOOL 1: list_dir
  Description: Lists all files and subdirectories at a given path.
  Usage: ACTION: list_dir | <path>

TOOL 2: read_file
  Description: Reads the text content of a single file.
  Usage: ACTION: read_file | <path>
  Note: You can natively read .txt, .md, .csv, .docx, .xlsx, .pdf, and .pptx files.

TOOL 3: copy_file
  Description: Copies a file from a source path to a destination path.
  Usage: ACTION: copy_file | <source> | <destination>

TOOL 4: write_file
  Description: Creates a new file or overwrites an existing file with the given content.
  Usage: ACTION: write_file | <path> | <content>
  Note: You can natively write to .docx and .xlsx files. For .xlsx, output raw CSV data (with commas).

TOOL 5: move_file
  Description: Moves or renames a file from a source path to a destination path.
  Usage: ACTION: move_file | <source> | <destination>

TOOL 6: search_in_files
  Description: Search for a specific text pattern inside all text files within a directory.
  Usage: ACTION: search_in_files | <directory> | <pattern>

TOOL 7: delete_file
  Description: Deletes a file. A backup is automatically saved before deletion.
  Usage: ACTION: delete_file | <path>

TOOL 8: create_dir
  Description: Creates a new empty directory.
  Usage: ACTION: create_dir | <path>

TOOL 9: append_to_file
  Description: Adds text to the END of an existing file (without overwriting it).
  Usage: ACTION: append_to_file | <path> | <content_to_add>

TOOL 10: file_info
  Description: Gets metadata about a file or folder (size, creation date, modification date).
  Usage: ACTION: file_info | <path>

TOOL 11: tree
  Description: Returns a recursive visual tree of the directory structure (up to 3 levels deep).
  Usage: ACTION: tree | <directory_path>
  Example: ACTION: tree | .
"""

TOOL_FUNCTIONS = {
    "list_dir": list_dir,
    "read_file": read_file,
    "copy_file": copy_file,
    "write_file": write_file,
    "move_file": move_file,
    "search_in_files": search_in_files,
    "delete_file": delete_file,
    "create_dir": create_dir,
    "append_to_file": append_to_file,
    "file_info": file_info,
    "tree": tree,
}
